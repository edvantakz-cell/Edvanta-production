from docx import Document
from pptx import Presentation
from pptx.util import Inches
import os
import uuid

BASE_PATH = "storage"

def save_lesson_doc(content):
    filename = f"{uuid.uuid4()}.docx"
    path = os.path.join(BASE_PATH, "lessons", filename)

    os.makedirs(os.path.dirname(path), exist_ok=True)

    doc = Document()
    doc.add_heading("Поурочный план", level=1)
    doc.add_paragraph(content)
    doc.save(path)

    return path

def save_presentation(content):
    filename = f"{uuid.uuid4()}.pptx"
    path = os.path.join(BASE_PATH, "presentations", filename)

    os.makedirs(os.path.dirname(path), exist_ok=True)

    prs = Presentation()
    slide_layout = prs.slide_layouts[1]

    for i in range(10):
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = f"Слайд {i+1}"
        slide.placeholders[1].text = content[:500]

    prs.save(path)
    return path